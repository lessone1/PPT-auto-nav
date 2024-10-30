from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 打开PPT文件
prs = Presentation("input.pptx")

# 获取每张幻灯片的第一个文本框作为标题
titles = []
for slide in prs.slides:
    slide_title = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text:
            # 直接读取第一个有文本的文本框作为标题
            slide_title = shape.text
            break
    titles.append(slide_title if slide_title else "（无标题）")

# 用户指定的导航栏分段方式，例如："2-5,7-12"
input_sections = input("请输入导航栏分段范围（例如：2-5,7-12）：")
# 解析用户输入为一个列表，格式为[(起始页, 终止页), ...]
sections = []
for part in input_sections.split(','):
    start, end = map(int, part.split('-'))
    sections.append((start, end))

# 遍历每个段落并生成导航栏
for start, end in sections:
    # 提取当前段的标题
    section_titles = titles[start - 1:end]
    ppt_width = prs.slide_width
    title_count = len(section_titles)
    nav_width = ppt_width / title_count  # 均分页面宽度

    # 为该段内的每张幻灯片生成导航栏
    for slide_idx in range(start - 1, end):
        slide = prs.slides[slide_idx]
        
        # 遍历当前段内的标题，为每个标题添加一个文本框
        for i, title in enumerate(section_titles):
            # 计算文本框位置，确保导航栏从页面最顶部开始，宽度正好填满页面
            left = i * nav_width  # 每个文本框左边距根据索引移动
            top = 0  # 设置为0，紧贴顶部
            width = nav_width
            height = Inches(0.5)  # 高度设置适中，确保不留空白

            # 创建导航栏文本框
            nav_shape = slide.shapes.add_textbox(left, top, width, height)
            nav_frame = nav_shape.text_frame
            nav_frame.word_wrap = True
            nav_frame.clear()  # 清除默认段落

            # 设置文本框的锚定模式为中间
            for paragraph in nav_frame.paragraphs:
                paragraph.space_before = Pt(0)  # 可以根据需要调整
                paragraph.space_after = Pt(0)    # 可以根据需要调整
                paragraph.alignment = PP_ALIGN.CENTER  # 文本水平居中

            # 设置文本框的锚定模式为垂直居中
            nav_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 设置当前页对应的文本框背景和字体
            if (start + i - 1) == slide_idx:
                # 当前页：蓝底白字
                fill = nav_shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(31, 78, 121)  # 当前页背景颜色 #1F4E79
                # nav_frame.text = f"{start + i}. {title}"  # 设置文本框内容
                nav_frame.text = f"{title}"  # 设置文本框内容
                p = nav_frame.paragraphs[0]
                p.font.size = Pt(12)  # 字体大小设置为12
                p.font.name = '微软雅黑'  # 字体样式设置为微软雅黑
                p.font.color.rgb = RGBColor(255, 255, 255)  # 字体颜色为白色
                # p.font.bold = True  # 字体加粗
                p.alignment = PP_ALIGN.CENTER
                
            else:
                # 其他页：背景颜色为 #BDD7EE，字体为黑色
                fill = nav_shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(189, 215, 238)  # 背景颜色 #BDD7EE
                nav_frame.text = f"{title}"  # 设置文本框内容
                p = nav_frame.paragraphs[0]
                p.font.size = Pt(12)  # 字体大小设置为12
                p.font.name = '微软雅黑'  # 字体样式设置为微软雅黑
                p.font.color.rgb = RGBColor(0, 0, 0)  # 字体颜色为黑色
                # p.font.bold = True  # 字体加粗
                p.alignment = PP_ALIGN.CENTER
                

            # 为文本框设置白色边框
            nav_shape.line.color.rgb = RGBColor(255, 255, 255)  # 边框颜色设为白色
            nav_shape.line.width = Pt(1)  # 边框宽度设为1pt

            # 添加超链接：点击文本框跳转到对应的幻灯片
            slide_number = start + i - 1  # 对应的幻灯片索引
            action = nav_shape.click_action
            action.hyperlink.address = f"#{slide_number + 1}"  # 超链接指向对应幻灯片

# 保存新PPT文件
prs.save("output.pptx")
